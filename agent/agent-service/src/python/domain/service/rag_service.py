import asyncio
import json
import time
from typing import Tuple, List, Dict, Any, Optional
from io import BytesIO
from datetime import timedelta
from urllib.parse import quote

from minio import Minio
from qdrant_client.http import models as rest

from infrastructure.common.logging.logging import logger
from infrastructure.common.error.errcode import ErrorCode
from infrastructure.config.sys_config import SysConfig
from infrastructure.repositories.rag_repository import RagRepository
from infrastructure.client.embedding_client import EmbeddingClient
from infrastructure.client.redis_client import RedisTemplete
from infrastructure.persistences.qdrant_persistence import QDrantPersistence
from infrastructure.client.llm.hsyq_client import HsyqClient
from infrastructure.client.neo4j_client import Neo4jClient
from interfaces.dto.rag_dto import (
    KnowledgeBaseCreateRequest,
    KnowledgeBaseUpdateRequest,
    KnowledgeBaseIdRequest,
    KnowledgeBaseListRequest,
    KnowledgeBaseInfo,
    KnowledgeBasePage,
    KnowledgeBaseStatus,
    KnowledgeBaseType,
    RagQueryRequest,
    RagQueryAnswer,
    FileVectorizationStatusRequest,
    FileIdWithKbRequest,
    FileUpdateStatusRequest,
    FileVectorizationStatusResponse,
    KnowledgeItemListRequest,
    KnowledgeItemInfo,
    KnowledgeSegmentListRequest,
    KnowledgeSegmentInfo,
    KnowledgeSegmentUpdateRequest,
    KnowledgeSegmentCreateRequest,
    KnowledgeSegmentIdRequest,
    OntologyNodeCreateRequest,
    OntologyNodeUpdateRequest,
    OntologyNodeDeleteRequest,
    HybridRagQueryRequest,
    HybridRagQueryResponse,
)


@logger()
class RagService:
    """
    RAG 领域服务：知识库元数据（MySQL）+ 文档向量（Qdrant）+ 知识嵌入（vllm-application Embedding）。
    """

    def __init__(self, config: SysConfig, mysql_client=None, qdrant_client=None, minio_client=None, neo4j_client=None):
        self.config = config
        self.repo = RagRepository(config, mysql_client)
        
        # 接收客户端实例
        self._mysql = mysql_client
        self._qdrant = qdrant_client
        self._minio_client = minio_client
        self._neo4j_client = neo4j_client

        # vllm-application Embedding 客户端（知识嵌入）
        embedding_base_url = config.get_vllm_embedding_base_url()
        embedding_api_key = config.get_vllm_embedding_api_key()
        self._embedding_client = (
            EmbeddingClient(embedding_base_url,timeout=1800, api_key=embedding_api_key)
            if embedding_base_url
            else None
        )
        if self._embedding_client and self._embedding_client.is_available():
            self.log.info("RagService: EmbeddingClient (vllm) enabled for knowledge embedding")

        # 支持的文件类型配置
        self._supported_file_types = config.get_supported_file_types()
        self.log.info(f"RagService: Supported file types: {sorted(self._supported_file_types)}")

        # Qdrant 用于存储 RAG 文档向量
        system_cfg = config.get_system_config() or {}
        persistence_cfg = system_cfg.get("persistence", {}) or {}
        qdrant_cfg = persistence_cfg.get("qdrant", {}) or {}
        self._vector_dim = int(qdrant_cfg.get("vector_dim", 512))
        self._rag_collection = qdrant_cfg.get("collection", "pg_agent_application_collection")
        
        # 使用传入的qdrant_client，如果没有则保持原有的初始化方式作为兼容
        if qdrant_client:
            self._qdrant = qdrant_client
            self.log.info(f"RagService: Using injected Qdrant client for collection {self._rag_collection}")
        else:
            self._qdrant: Optional[QDrantPersistence] = None
            try:
                host = qdrant_cfg.get("host", "127.0.0.1")
                port = int(qdrant_cfg.get("port", 6333))
                api_key = qdrant_cfg.get("api_key") or None
                self._qdrant = QDrantPersistence(
                    url=f"http://{host}:{port}",
                    api_key=api_key,
                    timeout=30.0,
                )
                self._qdrant.vector_size = self._vector_dim
                self._qdrant.init_collection(
                    self._rag_collection,
                    vectors_config=rest.VectorParams(
                        size=self._vector_dim,
                        distance=rest.Distance.COSINE,
                    ),
                )
                self.log.info(f"RagService: Qdrant RAG collection {self._rag_collection} (dim={self._vector_dim})")
            except Exception as e:
                self.log.warning(f"RagService: Qdrant init skipped: {e}")

        # 向量化进度使用 RedisTemplete（封装好的 Redis），供前端轮询 /file/status 获取实时进度
        self._rag_progress_prefix = "rag:file:progress:"
        self._rag_progress_ttl = 86400  # 24 小时
        if RedisTemplete.is_init:
            self.log.info("RagService: 向量化进度将写入 Redis（RedisTemplete），前端轮询 /file/status 可获取实时进度")
        else:
            self.log.info("RagService: RedisTemplete 未初始化，向量化进度仅通过 MySQL 更新，前端轮询仍可获取进度（按批更新）")

        # MinIO 用于读取已上传文件内容（用于 .txt 解析后做嵌入）
        minio_cfg = persistence_cfg.get("minio", {}) or {}
        self._minio_bucket = minio_cfg.get("bucket", "pg-agent-rag-files")
        
        # 使用传入的minio_client，如果没有则保持原有的初始化方式作为兼容
        if minio_client:
            self._minio_client = minio_client
            self.log.info(f"RagService: Using injected MinIO client for bucket {self._minio_bucket}")
        else:
            self._minio_client: Optional[Minio] = None
            try:
                endpoint = minio_cfg.get("endpoint", "127.0.0.1:9000")
                self._minio_client = Minio(
                    endpoint,
                    access_key=minio_cfg.get("access_key", "admin"),
                    secret_key=minio_cfg.get("secret_key", "admin123"),
                    secure=bool(minio_cfg.get("secure", False)),
                )
                self.log.info(f"RagService: MinIO client initialized for bucket {self._minio_bucket}")
            except Exception as e:
                self.log.warning(f"RagService: MinIO client init skipped: {e}")

        # Neo4j 客户端（可选，用于本体提取）
        if neo4j_client:
            self._neo4j_client = neo4j_client
            self.log.info("RagService: Using injected Neo4j client")
        else:
            self._neo4j_client = None
            try:
                # 尝试初始化 Neo4j 客户端
                # 注意：Neo4jClient 内部会读取配置，若连接失败会记录错误但不会抛出异常
                client = Neo4jClient(config)
                self._neo4j_client = client
                if client._driver:
                    self.log.info("RagService: Neo4j client initialized successfully")
                else:
                    self.log.warning("RagService: Neo4j client initialized but driver is None (will retry on first use)")
            except Exception as e:
                self.log.warning(f"RagService: Neo4j client init skipped: {e}")
                self._neo4j_client = None

        # LLM 客户端（用于本体提取等）
        self._llm_client = HsyqClient(config)
        self.log.info(f"RagService: HsyqClient initialized with model {self._llm_client.model}")

        self.log.info("RagService initialized with RagRepository")

    # ============== 工具方法 ==============

    def _map_status_from_db(self, status_val: int) -> KnowledgeBaseStatus:
        """
        将 rag_tbl.status (0/1) 映射为枚举：
        - 1 -> ENABLED
        - 0 -> DISABLED
        """
        return KnowledgeBaseStatus.ENABLED if status_val == 1 else KnowledgeBaseStatus.DISABLED

    def _map_status_to_db(self, status: KnowledgeBaseStatus) -> int:
        """
        将枚举映射为 rag_tbl.status：
        - ENABLED / BUILDING -> 1
        - DISABLED -> 0
        """
        if status == KnowledgeBaseStatus.DISABLED:
            return 0
        return 1

    def _map_type_from_db(self, type_str: str) -> KnowledgeBaseType:
        try:
            return KnowledgeBaseType(type_str)
        except Exception:
            return KnowledgeBaseType.DOCUMENT

    def _row_to_kb(self, row: Dict[str, Any]) -> KnowledgeBaseInfo:
        """将 rag_tbl 行记录映射为 DTO"""
        create_time = row.get("create_time")
        update_time = row.get("update_time")
        created_at_ms = (
            int(create_time.timestamp() * 1000) if create_time is not None else 0
        )
        updated_at_ms = (
            int(update_time.timestamp() * 1000) if update_time is not None else 0
        )

        return KnowledgeBaseInfo(
            id=row.get("id"),
            name=row.get("rag_name", ""),
            kb_type=self._map_type_from_db(row.get("type", "document")),
            status=self._map_status_from_db(int(row.get("status", 1))),
            created_at=created_at_ms,
            updated_at=updated_at_ms,
            file_capacity=int(row.get("file_capacity", 0) or 0),
            file_count=int(row.get("file_count", 0) or 0),
            binding_user=int(row.get("binding_user_id", 0) or 0),
        )

    # ============== 知识嵌入（vllm-application Embedding + Qdrant） ==============

    def _rag_progress_key(self, file_id: int) -> str:
        return f"{self._rag_progress_prefix}{file_id}"

    def set_rag_file_progress(
        self,
        file_id: int,
        embedding_process: int,
        embedding_progress: float = 0.0,
        file_char_count: int = 0,
        file_embedding_offset: int = 0,
        error_message: Optional[str] = None,
    ) -> None:
        """写入 Redis（RedisTemplete）：当前文件的向量化进度，供前端轮询。value 中含 file_id 便于 Redis 客户端识别编号。"""
        if not RedisTemplete.is_init:
            return
        try:
            data = {
                "file_id": file_id,
                "embedding_process": embedding_process,
                "embedding_progress": embedding_progress,
                "file_char_count": file_char_count,
                "file_embedding_offset": file_embedding_offset,
            }
            if error_message:
                data["error_message"] = error_message
            key = self._rag_progress_key(file_id)
            RedisTemplete.set(key, json.dumps(data), self._rag_progress_ttl)
        except Exception as e:
            self.log.warning(f"set_rag_file_progress failed: {e}")

    def clear_rag_file_progress(self, file_id: int) -> None:
        """文件完成后清理对应的 Redis 进度信息。"""
        if not RedisTemplete.is_init:
            self.log.warning(f"clear rag file process progress failure: RedisTemplete not initialized")
            return
        try:
            key = self._rag_progress_key(file_id)
            RedisTemplete.delete(key)
        except Exception as e:
            self.log.warning(f"clear_rag_file_progress failed: {e}")

    def cleanup_file_resources(
        self,
        kb_id: int,
        file_id: int,
        object_name: Optional[str] = None,
        bucket: Optional[str] = None,
    ) -> None:
        """
        清理单个文件相关资源：
        - Redis 进度
        - Qdrant 向量（按 file_id）
        - MinIO 对象
        - MySQL rag_file_tbl 记录
        """
        # Redis
        self.clear_rag_file_progress(file_id)

        # Qdrant 向量
        if self._qdrant:
            try:
                filter_condition = rest.Filter(
                    must=[
                        rest.FieldCondition(
                            key="file_id",
                            match=rest.MatchAny(any=[file_id]),
                        )
                    ]
                )
                err_del, n_pts = self._qdrant.delete(
                    filter=filter_condition,
                    collection_name=self._rag_collection,
                )
                if err_del != ErrorCode.SUCCESS:
                    self.log.warning(f"cleanup_file_resources: Qdrant delete failed err={err_del} for kb_id={kb_id}, file_id={file_id}")
                else:
                    self.log.info(f"cleanup_file_resources: Qdrant deleted {n_pts} points for kb_id={kb_id}, file_id={file_id}")
            except Exception as e:
                self.log.warning(f"cleanup_file_resources: Qdrant delete exception for kb_id={kb_id}, file_id={file_id}: {e}")

        # Neo4j 本体
        self.delete_ontology_by_file(file_id)

        # MinIO 对象
        obj_name = object_name
        if not obj_name:
            try:
                err_row, row = self.repo.get_file_by_id(file_id)
                if err_row == ErrorCode.SUCCESS and row and row.get("file_path"):
                    obj_name = row.get("file_path")
            except Exception:
                obj_name = None
        if obj_name and self._minio_client:
            try:
                b = bucket or self._minio_bucket
                self._minio_client.remove_object(b, obj_name)
                self.log.info(f"cleanup_file_resources: MinIO removed object '{obj_name}' for kb_id={kb_id}, file_id={file_id}")
            except Exception as e:
                self.log.warning(f"cleanup_file_resources: MinIO remove_object exception for '{obj_name}': {e}")

        # MySQL 知识片段记录
        try:
            err_del_seg, _ = self.repo.delete_segments_by_file(file_id)
            if err_del_seg != ErrorCode.SUCCESS:
                self.log.warning(f"cleanup_file_resources: delete_segments_by_file failed err={err_del_seg} for file_id={file_id}")
        except Exception as e:
            self.log.warning(f"cleanup_file_resources: delete_segments_by_file exception for file_id={file_id}: {e}")

        # MySQL 文件记录
        try:
            err_del_file, _ = self.repo.delete_file_by_id(file_id)
            if err_del_file != ErrorCode.SUCCESS:
                self.log.warning(f"cleanup_file_resources: delete_file_by_id failed err={err_del_file} for kb_id={kb_id}, file_id={file_id}")
        except Exception as e:
            self.log.warning(f"cleanup_file_resources: delete_file_by_id exception for kb_id={kb_id}, file_id={file_id}: {e}")

    def get_rag_file_progress(self, file_id: int) -> Optional[Dict[str, Any]]:
        """从 Redis（RedisTemplete）读取文件向量化进度；无则返回 None。"""
        if not RedisTemplete.is_init:
            return None
        try:
            key = self._rag_progress_key(file_id)
            val = RedisTemplete.get(key)
            if val:
                return json.loads(val) if isinstance(val, str) else json.loads(val.decode("utf-8"))
        except Exception as e:
            self.log.warning(f"get_rag_file_progress failed: {e}")
        return None

    def embed_texts(
        self,
        texts: List[str],
        dimensions: Optional[int] = None,
    ) -> Tuple[ErrorCode, List[List[float]]]:
        """
        调用 vllm-application 将文本列表转为向量。
        dimensions 不传时使用配置的 vector_dim（与 Qdrant 一致）。
        """
        if not self._embedding_client or not self._embedding_client.is_available():
            self.log.warning("embed_texts: EmbeddingClient not available")
            return ErrorCode.SERVICE_UNAVAILABLE, []
        dim = dimensions if dimensions is not None else self._vector_dim
        return self._embedding_client.embed(texts, dimensions=dim)

    @staticmethod
    def _chunk_text(text: str, chunk_size: int = 500, overlap: int = 50) -> List[str]:
        """简单按字符数切分，overlap 为重叠长度。"""
        if not text or chunk_size <= 0:
            return []
        text = text.strip()
        if not text:
            return []
        step = max(1, chunk_size - overlap)
        chunks = []
        for i in range(0, len(text), step):
            chunk = text[i : i + chunk_size]
            if chunk:
                chunks.append(chunk)
        return chunks

    def index_chunks_for_file(
        self,
        kb_id: int,
        file_id: int,
        binding_user_id: int,
        chunks: List[str],
        filename: str = "",
        file_char_count: int = 0,
        batch_size: int = 20,
    ) -> ErrorCode:
        """
        将文档分块向量化并写入 Qdrant。若 file_char_count > 0 则按批处理并写入 Redis 进度。
        payload: rag_id, file_id, chunk_index, text, binding_user_id
        """
        if not chunks:
            return ErrorCode.SUCCESS
        if not self._qdrant:
            self.log.warning("index_chunks_for_file: Qdrant not available")
            return ErrorCode.SERVICE_UNAVAILABLE
        total_chars = sum(len(c) for c in chunks)
        use_progress = file_char_count > 0 and RedisTemplete.is_init
        offset_chars = 0
        for start in range(0, len(chunks), batch_size):
            batch = chunks[start : start + batch_size]
            err, embeddings = self.embed_texts(batch)
            if err != ErrorCode.SUCCESS or len(embeddings) != len(batch):
                self.log.error(f"index_chunks_for_file: embed failed err={err}, got {len(embeddings)} vectors")
                return err if err != ErrorCode.SUCCESS else ErrorCode.DATA_FORMAT_INVALID
            points = []
            for j, (vec, text) in enumerate(zip(embeddings, batch)):
                i = start + j
                if len(vec) != self._vector_dim:
                    self.log.warning(f"Chunk {i} vector dim {len(vec)} != {self._vector_dim}, skip")
                    continue

                # Qdrant 仅接受 unsigned int 或 UUID，使用确定性整数 ID 避免冲突
                point_id = (kb_id << 40) | (file_id << 20) | i

                # 插入 MySQL 知识片段
                seg_data = {
                    "file_id": file_id,
                    "binding_user_id": binding_user_id,
                    "title": f"{filename}#{i+1}" if filename else f"Segment#{i+1}",
                    "content": text,
                    "summary": text[:200],
                    "sort_index": i,
                    "status": 2,  # 启用
                    "snippet_type": "text",  # 文本
                    "recall_count": 0,
                    "qdrant_vector_id": str(point_id),
                }
                err_ins, seg_id = self.repo.insert_segment(seg_data)
                if err_ins != ErrorCode.SUCCESS:
                    self.log.error(f"index_chunks_for_file: insert_segment failed err={err_ins}")
                    return err_ins

                points.append({
                    "id": point_id,
                    "vector": vec,
                    "payload": {
                        "rag_id": kb_id,
                        "file_id": file_id,
                        "segment_id": seg_id,
                        "chunk_index": i,
                        "text": text[:2000],
                        "binding_user_id": binding_user_id,
                    },
                })
            err, n = self._qdrant.insert(points)
            if err != ErrorCode.SUCCESS:
                self.log.error(f"index_chunks_for_file: Qdrant insert failed err={err}")
                return err
            offset_chars += sum(len(c) for c in batch)
            # 每批都更新 MySQL 的 file_embedding_offset，前端轮询 /file/status 时才能看到进度（无论是否启用 Redis）
            self.repo.update_file_record(file_id, {"file_embedding_offset": offset_chars})
            if use_progress and total_chars > 0:
                progress = min(offset_chars / file_char_count, 1.0)
                self.set_rag_file_progress(
                    file_id,
                    embedding_process=1,
                    embedding_progress=progress,
                    file_char_count=file_char_count,
                    file_embedding_offset=offset_chars,
                )
        # 全部写入 Qdrant 后，在 Redis 标记完成，供 status 只查 Redis
        if use_progress and file_char_count > 0:
            self.set_rag_file_progress(
                file_id,
                embedding_process=3,
                embedding_progress=1.0,
                file_char_count=file_char_count,
                file_embedding_offset=offset_chars,
            )
        self.log.info(f"index_chunks_for_file: kb_id={kb_id} file_id={file_id} indexed {len(chunks)} chunks")
        return ErrorCode.SUCCESS

    def _extract_and_save_ontology(self, text: str, file_id: Optional[int] = None, strict: bool = False):
        """
        利用 LLM 提取本体并写入 Neo4j
        strict: If True, raise exception on failure (for Saga rollback)
        """
        if not self._neo4j_client or not self._llm_client:
            return

        # 从配置中读取 Prompt 模板
        system_cfg = self.config.get_system_config() or {}
        ontology_cfg = system_cfg.get("neo4j_ontology", {}) or {}
        
        # 获取最大长度配置，默认为 0 (不限制)，允许使用配置的嵌入文本提取本体
        try:
            max_length = int(ontology_cfg.get("max_length", 0))
        except (ValueError, TypeError):
            max_length = 0

        if max_length > 0:
            content_snippet = text[:max_length]
        else:
            content_snippet = text

        prompt_template = ontology_cfg.get("prompt")
        
        if not prompt_template:
            # 默认模板（如果配置缺失）
            prompt_template = """
你是一个知识图谱专家。请从以下文本中提取关键实体（Entity）和关系（Relationship），构建本体结构。
请忽略无关紧要的词汇，只提取核心概念。

文本内容：
{text}

请严格按照以下 JSON 格式返回结果（不要包含 Markdown 代码块标记）：
{
    "nodes": [
        { "name": "实体名称", "label": "实体类型", "properties": { "desc": "简短描述" } }
    ],
    "relationships": [
        { "head": "头实体名称", "tail": "尾实体名称", "type": "关系类型" }
    ]
}
"""
        
        # 替换占位符
        prompt = prompt_template.replace("{text}", content_snippet)
        try:
            self.log.debug(f"_extract_and_save_ontology: prompt={prompt}")
            resp = self._llm_client.chat([{"role": "user", "content": prompt}])
            if not resp:
                self.log.debug("_extract_and_save_ontology: LLM response is empty")
                return
            self.log.debug(f"_extract_and_save_ontology: LLM response={resp}")
                
            # 清理可能的 Markdown 标记
            clean_resp = resp.strip()
            if clean_resp.startswith("```json"):
                clean_resp = clean_resp[7:]
            if clean_resp.startswith("```"):
                clean_resp = clean_resp[3:]
            if clean_resp.endswith("```"):
                clean_resp = clean_resp[:-3]
            
            data = json.loads(clean_resp)
            nodes = data.get("nodes", [])
            rels = data.get("relationships", [])
            
            self.log.info(f"Extracted {len(nodes)} nodes and {len(rels)} relationships from text")
            
            # 写入 Neo4j
            for node in nodes:
                name = node.get("name")
                label = node.get("label", "Entity")
                props = node.get("properties", {})
                if not name: continue
                
                # 简单的防注入处理（label 不能参数化）
                safe_label = "".join(c for c in label if c.isalnum() or c == '_')
                if not safe_label: safe_label = "Entity"
                
                cypher = f"MERGE (n:`{safe_label}` {{name: $name}}) SET n += $props"
                self._neo4j_client.execute_query(cypher, {"name": name, "props": props})
                
            for rel in rels:
                head = rel.get("head")
                tail = rel.get("tail")
                rtype = rel.get("type", "RELATED_TO")
                if not head or not tail: continue
                
                safe_rtype = "".join(c for c in rtype if c.isalnum() or c == '_').upper()
                if not safe_rtype: safe_rtype = "RELATED_TO"
                
                # 仅通过 name 匹配（假设 name 全局唯一或能接受模糊匹配）
                cypher = f"""
                MATCH (h {{name: $head}})
                MATCH (t {{name: $tail}})
                MERGE (h)-[r:`{safe_rtype}`]->(t)
                SET r.file_id = $file_id
                """
                self._neo4j_client.execute_query(cypher, {"head": head, "tail": tail, "file_id": file_id})
                
        except Exception as e:
            self.log.warning(f"Ontology extraction failed: {e}")
            if strict:
                raise e

    def delete_ontology_by_file(self, file_id: int):
        """
        Rollback ontology: delete relationships created by this file
        """
        if not self._neo4j_client:
            return
        
        try:
            query = "MATCH ()-[r]-() WHERE r.file_id = $file_id DELETE r"
            self._neo4j_client.execute_query(query, {"file_id": file_id})
            self.log.info(f"Deleted ontology relationships for file_id={file_id}")
        except Exception as e:
            self.log.warning(f"Failed to delete ontology for file_id={file_id}: {e}")

    def embed_and_index_file(
        self,
        kb_id: int,
        file_id: int,
        binding_user_id: int,
        object_name: str,
        bucket: Optional[str] = None,
        chunk_size: int = 500,
    ) -> ErrorCode:
        """
        从 MinIO 读取文件内容，支持 PDF、TXT、Office、HTML；分块后调用 vllm 嵌入并写入 Qdrant。
        若文件类型不支持或读取/解析失败，则视为该文件向量化失败。
        """
        if not object_name or not self._minio_client:
            self.log.warning("embed_and_index_file: missing object_name or MinIO client, treat as failure")
            return ErrorCode.OPERATION_FAILED
        
        # 检查文件扩展名
        file_ext = object_name.lower().split('.')[-1] if '.' in object_name.lower() else ''
        
        if file_ext not in self._supported_file_types:
            self.log.warning(f"embed_and_index_file: unsupported file type '{file_ext}', treat as failure")
            return ErrorCode.NOT_SUPPORTED
        
        b = bucket or self._minio_bucket
        try:
            resp = self._minio_client.get_object(b, object_name)
            content = resp.read()
            resp.close()
        except Exception as e:
            self.log.warning(f"embed_and_index_file: MinIO get_object failed: {e}")
            return ErrorCode.FILE_NOT_FOUND
        
        # 根据文件类型提取文本
        text = self._extract_text_from_file(content, file_ext)
        if not text:
            self.log.warning(f"embed_and_index_file: failed to extract text from {file_ext} file")
            return ErrorCode.INVALID_DATA
        
        chunks = self._chunk_text(text, chunk_size=chunk_size)
        if not chunks:
            self.log.warning("embed_and_index_file: chunk_text returned empty chunks, treat as failure")
            return ErrorCode.INVALID_DATA
        file_char_count = len(text)
        # 写回实际字符数，便于 /file/status 计算进度（PDF 等插入时 char_count 常为 0）
        self.repo.update_file_record(file_id, {"file_char_count": file_char_count})
        self.set_rag_file_progress(
            file_id,
            embedding_process=1,
            embedding_progress=0.0,
            file_char_count=file_char_count,
            file_embedding_offset=0,
        )
        err = self.index_chunks_for_file(
            kb_id, file_id, binding_user_id, chunks,
            filename=object_name,
            file_char_count=file_char_count,
        )
        if err != ErrorCode.SUCCESS:
            return err

        try:
            # 向量化成功后，标记文件已完成嵌入
            self.repo.update_file_record(file_id, {"embedding_process": 3})

            # 重新统计知识库下已完成文件的容量与数量，并写回 rag_tbl
            stats_err, (total_chars, file_cnt) = self.repo.calc_kb_file_stats(kb_id)
            if stats_err == ErrorCode.SUCCESS:
                upd_err, _ = self.repo.update_kb_file_stats(kb_id, total_chars, file_cnt)
                if upd_err != ErrorCode.SUCCESS:
                    self.log.warning(f"update_kb_file_stats failed: {upd_err} for kb_id={kb_id}")
            else:
                self.log.warning(f"calc_kb_file_stats failed: {stats_err} for kb_id={kb_id}")
            
            # 自动提取本体
            self._extract_and_save_ontology(text, file_id=file_id)

        except Exception as e:
            self.log.warning(f"update kb file stats exception for kb_id={kb_id}: {e}")

        return ErrorCode.SUCCESS

    def create_knowledge_from_file_content_saga(
        self,
        kb_id: int,
        filename: str,
        content: bytes,
        binding_user_id: int,
        chunk_size: int = 500
    ) -> ErrorCode:
        """
        Create knowledge base from file content using Saga pattern.
        Flow: Cut -> MySQL -> Qdrant -> Neo4j -> MinIO -> Relations
        """
        # 1. Text Extraction & Chunking
        file_ext = filename.lower().split('.')[-1] if '.' in filename.lower() else ''
        text = self._extract_text_from_file(content, file_ext)
        if not text:
            return ErrorCode.INVALID_DATA
            
        chunks = self._chunk_text(text, chunk_size=chunk_size)
        if not chunks:
            return ErrorCode.INVALID_DATA
            
        file_char_count = len(text)
        
        # 2. Saga Step 1: MySQL (File Record) - "Content parent index"
        file_data = {
            "rag_id": kb_id,
            "file_name": filename,
            "file_path": "", # Placeholder, will update later
            "extension": file_ext,
            "status": 1, # Enabled
            "file_char_count": file_char_count,
            "embedding_process": 1, # Processing
            "binding_user_id": binding_user_id,
            "upload_time": time.strftime("%Y-%m-%d %H:%M:%S"),
        }
        err, file_id = self.repo.insert_file_record(file_data)
        if err != ErrorCode.SUCCESS:
            return err
            
        uploaded_object_name = None

        # Define compensation (Rollback) function
        def rollback():
            self.log.warning(f"Saga Rollback initiated for file_id={file_id}")
            self.cleanup_file_resources(kb_id, file_id, object_name=uploaded_object_name)
            
        try:
            # 3. Saga Step 2: MySQL (Segments) & Qdrant (Vectors)
            err = self.index_chunks_for_file(
                kb_id, file_id, binding_user_id, chunks,
                filename=filename,
                file_char_count=file_char_count
            )
            if err != ErrorCode.SUCCESS:
                raise Exception(f"Index chunks failed: {err}")
                
            # 4. Saga Step 3: Neo4j (Ontology)
            # Use strict=True to enforce Saga
            self._extract_and_save_ontology(text, file_id=file_id, strict=True)
            
            # 5. Saga Step 4: MinIO (Upload)
            object_name = f"{kb_id}/{file_id}/{filename}"
            if self._minio_client:
                self._minio_client.put_object(
                    self._minio_bucket,
                    object_name,
                    BytesIO(content),
                    len(content)
                )
                uploaded_object_name = object_name
            else:
                 raise Exception("MinIO client not available")
                 
            # 6. Saga Step 5: Relations (Update MySQL with MinIO path)
            update_data = {
                "file_path": object_name,
                "embedding_process": 3 # Completed
            }
            err, _ = self.repo.update_file_record(file_id, update_data)
            if err != ErrorCode.SUCCESS:
                raise Exception(f"Update file record failed: {err}")
                
            # Update KB stats (Best effort)
            self.repo.calc_kb_file_stats(kb_id)
            
        except Exception as e:
            self.log.error(f"Saga failed: {e}")
            rollback()
            return ErrorCode.OPERATION_FAILED
            
        return ErrorCode.SUCCESS

    def _extract_text_from_file(self, content: bytes, file_ext: str) -> str:
        """
        根据文件类型提取文本内容
        """
        import io
        import re
        
        try:
            if file_ext == 'txt':
                return content.decode("utf-8", errors="replace")
            
            elif file_ext in ['html', 'htm']:
                # HTML文件文本提取
                try:
                    from bs4 import BeautifulSoup
                    soup = BeautifulSoup(content, 'html.parser')
                    # 移除script和style标签
                    for script in soup(["script", "style"]):
                        script.decompose()
                    return soup.get_text(separator=' ', strip=True)
                except ImportError:
                    self.log.warning("BeautifulSoup not available, falling back to regex for HTML")
                    # 简单的HTML标签移除
                    text = re.sub(r'<[^>]+>', ' ', content.decode("utf-8", errors="replace"))
                    return re.sub(r'\s+', ' ', text).strip()
            
            elif file_ext == 'pdf':
                # PDF文件文本提取
                try:
                    import PyPDF2
                    pdf_reader = PyPDF2.PdfReader(io.BytesIO(content))
                    text = ""
                    for page in pdf_reader.pages:
                        text += page.extract_text() + "\n"
                    return text.strip()
                except ImportError:
                    self.log.warning("PyPDF2 not available, cannot extract text from PDF")
                    return ""
            
            elif file_ext in ['doc', 'docx']:
                # Word文档文本提取
                try:
                    if file_ext == 'docx':
                        import docx
                        doc = docx.Document(io.BytesIO(content))
                        text = ""
                        for paragraph in doc.paragraphs:
                            text += paragraph.text + "\n"
                        return text.strip()
                    else:
                        # .doc文件需要python-docx2txt
                        import docx2txt
                        return docx2txt.process(io.BytesIO(content))
                except ImportError:
                    self.log.warning("python-docx or docx2txt not available, cannot extract text from Word documents")
                    return ""
            
            elif file_ext in ['xls', 'xlsx']:
                # Excel文件文本提取
                try:
                    import pandas as pd
                    df = pd.read_excel(io.BytesIO(content))
                    # 将所有单元格内容转换为文本
                    text = ""
                    for _, row in df.iterrows():
                        row_text = " ".join([str(cell) for cell in row if pd.notna(cell)])
                        text += row_text + "\n"
                    return text.strip()
                except ImportError:
                    self.log.warning("pandas not available, cannot extract text from Excel files")
                    return ""
            
            elif file_ext in ['ppt', 'pptx']:
                # PowerPoint文件文本提取
                try:
                    if file_ext == 'pptx':
                        import pptx
                        prs = pptx.Presentation(io.BytesIO(content))
                        text = ""
                        for slide in prs.slides:
                            for shape in slide.shapes:
                                if hasattr(shape, "text"):
                                    text += shape.text + "\n"
                        return text.strip()
                    else:
                        # .ppt文件支持较少，可能需要其他库
                        self.log.warning(".ppt files not fully supported, consider converting to .pptx")
                        return ""
                except ImportError:
                    self.log.warning("python-pptx not available, cannot extract text from PowerPoint files")
                    return ""
            
            else:
                self.log.warning(f"Unsupported file type: {file_ext}")
                return ""
                
        except Exception as e:
            self.log.error(f"Error extracting text from {file_ext} file: {e}")
            return ""

    def check_file_vectorized(self, file_name: str, kb_id: int) -> Tuple[ErrorCode, bool]:
        """
        检查文件是否已经向量化
        先检查MySQL中是否存在文件记录，再检查Qdrant中是否存在向量化数据
        """
        # 1. 检查MySQL中是否存在文件记录
        err, file_record = self.repo.get_file_by_name_and_kb(file_name, kb_id)
        if err != ErrorCode.SUCCESS:
            return err, False
        
        if not file_record:
            # MySQL中不存在文件记录，可以插入
            return ErrorCode.SUCCESS, False
        
        # 2. 检查Qdrant中是否存在向量化数据
        file_id = file_record.get("id")
        if not self._qdrant:
            self.log.warning("check_file_vectorized: Qdrant not available")
            return ErrorCode.SERVICE_UNAVAILABLE, False
        
        err, qdrant_records = self._qdrant.query_by_file_id(
            file_id=file_id,
            collection_name=self._rag_collection,
            limit=1
        )
        
        if err != ErrorCode.SUCCESS:
            return err, False
        
        is_vectorized = len(qdrant_records) > 0
        return ErrorCode.SUCCESS, is_vectorized

    # ============== 知识库管理 ==============

    async def create_knowledge_base(
        self, req: KnowledgeBaseCreateRequest, user_id: int
    ) -> Tuple[ErrorCode, object]:
        """创建知识库 -> rag_tbl，绑定当前用户 binding_user_id"""
        err, exists = await asyncio.to_thread(
            self.repo.exists_by_name, req.name, binding_user_id=user_id
        )
        if err != ErrorCode.SUCCESS:
            self.log.error(f"exists check failed when creating kb: {err}")
            return err, "数据库错误"
        if exists:
            msg = f"知识库名称已存在: {req.name}"
            self.log.warning(msg)
            return ErrorCode.RESOURCE_ALREADY_EXISTS, msg

        data = {
            "rag_name": req.name,
            "type": req.kb_type.value,
            "status": self._map_status_to_db(KnowledgeBaseStatus.ENABLED),
            "file_capacity": 0,
            "file_count": 0,
            "binding_user_id": user_id,
        }
        err, kb_id = await asyncio.to_thread(self.repo.insert_knowledge_base, data)
        if err != ErrorCode.SUCCESS:
            self.log.error(f"insert rag_tbl failed: {err}")
            return err, "创建知识库失败"

        # 读取刚插入的记录，构造 DTO
        err, row = await asyncio.to_thread(self.repo.get_knowledge_base_by_id, kb_id)
        if err != ErrorCode.SUCCESS or not row:
            self.log.error(f"select rag_tbl after insert failed: {err}")
            return err or ErrorCode.DATA_NOT_FOUND, "读取知识库失败"

        kb = self._row_to_kb(row)
        self.log.info(f"Created knowledge base in DB: id={kb.id}, name={kb.name}")
        return ErrorCode.SUCCESS, kb

    async def update_knowledge_base(
        self, req: KnowledgeBaseUpdateRequest, user_id: int
    ) -> Tuple[ErrorCode, object]:
        """更新知识库信息 -> rag_tbl（仅允许操作当前用户绑定的知识库）"""
        err, row = await asyncio.to_thread(
            self.repo.get_knowledge_base_by_id, req.id, binding_user_id=user_id
        )
        if err != ErrorCode.SUCCESS or not row:
            msg = f"知识库不存在或无权操作: {req.id}"
            self.log.warning(msg)
            return ErrorCode.DATA_NOT_FOUND, msg

        if req.name and req.name != row.get("rag_name"):
            err, exists = await asyncio.to_thread(
                self.repo.exists_by_name, req.name, binding_user_id=user_id, exclude_id=req.id
            )
            if err != ErrorCode.SUCCESS:
                self.log.error(f"exists check failed when updating kb: {err}")
                return err, "数据库错误"
            if exists:
                msg = f"知识库名称已存在: {req.name}"
                self.log.warning(msg)
                return ErrorCode.RESOURCE_ALREADY_EXISTS, msg

        update_data: Dict[str, Any] = {}
        if req.name:
            update_data["rag_name"] = req.name
        if req.status is not None:
            update_data["status"] = self._map_status_to_db(req.status)

        if not update_data:
            # 没有任何需要更新的字段，返回当前数据
            kb = self._row_to_kb(row)
            return ErrorCode.SUCCESS, kb

        err, _ = await asyncio.to_thread(
            self.repo.update_knowledge_base, req.id, update_data
        )
        if err != ErrorCode.SUCCESS:
            self.log.error(f"update rag_tbl failed: {err}, data={update_data}")
            return err, "更新知识库失败"

        # 返回更新后的数据
        err, row2 = await asyncio.to_thread(self.repo.get_knowledge_base_by_id, req.id)
        if err != ErrorCode.SUCCESS or not row2:
            self.log.error(f"select rag_tbl after update failed: {err}")
            return err or ErrorCode.DATA_NOT_FOUND, "读取知识库失败"

        kb = self._row_to_kb(row2)
        self.log.info(f"Updated knowledge base in DB: id={kb.id}")
        return ErrorCode.SUCCESS, kb

    async def delete_knowledge_base(
        self, req: KnowledgeBaseIdRequest, user_id: int
    ) -> Tuple[ErrorCode, object]:
        """
        删除知识库（仅允许删除当前用户绑定的知识库），并级联删除：
        - Qdrant 中该知识库下所有文件的向量
        - MinIO 中该知识库下所有文档对象
        - MySQL rag_file_tbl 中该知识库下所有文件记录
        - MySQL rag_tbl 中知识库主记录
        使用事务保证 MySQL 内「先删文件表、再删知识库表」的原子性。
        """
        kb_id = req.id
        err, row = await asyncio.to_thread(
            self.repo.get_knowledge_base_by_id, kb_id, binding_user_id=user_id
        )
        if err != ErrorCode.SUCCESS or not row:
            msg = f"知识库不存在或无权操作: {kb_id}"
            self.log.warning(msg)
            return ErrorCode.DATA_NOT_FOUND, msg

        # 1. 查出该知识库下所有文件（用于删 Qdrant、MinIO）
        err, file_rows = await asyncio.to_thread(
            self.repo.list_files_by_kb, kb_id, limit=50000
        )
        if err != ErrorCode.SUCCESS:
            self.log.error(f"list_files_by_kb failed: {err}, kb_id={kb_id}")
            return err, "查询知识库文件列表失败"
        file_ids = [r["id"] for r in file_rows]
        file_paths = [r.get("file_path") for r in file_rows if r.get("file_path")]

        # 2. 删除 Qdrant 中该知识库下所有 file_id 的向量
        qdrant_deleted = 0
        if file_ids and self._qdrant:
            try:
                filter_condition = rest.Filter(
                    must=[
                        rest.FieldCondition(
                            key="file_id",
                            match=rest.MatchAny(any=file_ids),
                        )
                    ]
                )
                err_del, n_pts = await asyncio.to_thread(
                    self._qdrant.delete,
                    filter=filter_condition,
                    collection_name=self._rag_collection,
                )
                if err_del == ErrorCode.SUCCESS:
                    qdrant_deleted = n_pts
                    self.log.info(f"Qdrant: deleted points for kb_id={kb_id}, file_ids={len(file_ids)}, points={qdrant_deleted}")
                else:
                    self.log.warning(f"Qdrant delete by file_ids failed: {err_del}, kb_id={kb_id}")
            except Exception as e:
                self.log.warning(f"Qdrant delete exception: {e}, kb_id={kb_id}")

        # 3. 删除 MinIO 中该知识库下所有文档对象（放入线程避免阻塞）
        minio_deleted = 0
        if file_paths and self._minio_client:
            def _remove_minio_objects(bucket: str, paths: list) -> int:
                cnt = 0
                for obj_name in paths:
                    try:
                        self._minio_client.remove_object(bucket, obj_name)
                        cnt += 1
                    except Exception:
                        pass
                return cnt
            minio_deleted = await asyncio.to_thread(
                _remove_minio_objects, self._minio_bucket, file_paths
            )
            if minio_deleted > 0:
                self.log.info(f"MinIO: removed {minio_deleted} objects for kb_id={kb_id}")

        # 4. MySQL 事务：先删文件表，再删知识库表（整段事务放入线程，避免阻塞事件循环）
        def _do_delete_tx():
            if self.repo.begin_transaction() != ErrorCode.SUCCESS:
                return ErrorCode.DATABASE_TRANSACTION_ERROR, "开启事务失败", 0
            try:
                # 0. 删除知识片段
                err0, _ = self.repo.delete_segments_by_kb(kb_id)
                if err0 != ErrorCode.SUCCESS:
                    self.repo.rollback()
                    return err0, "删除知识片段失败", 0

                err1, n_files = self.repo.delete_files_by_kb(kb_id)
                if err1 != ErrorCode.SUCCESS:
                    self.repo.rollback()
                    return err1, "删除知识库文件记录失败", 0
                err2, n_kb = self.repo.delete_knowledge_base(kb_id)
                if err2 != ErrorCode.SUCCESS:
                    self.repo.rollback()
                    return err2, "删除知识库失败", 0
                if self.repo.commit() != ErrorCode.SUCCESS:
                    self.repo.rollback()
                    return ErrorCode.DATABASE_TRANSACTION_ERROR, "提交事务失败", 0
                return ErrorCode.SUCCESS, None, n_files
            except Exception as e:
                self.repo.rollback()
                return ErrorCode.DATABASE_TRANSACTION_ERROR, f"删除知识库事务异常: {e}", 0

        tx_err, tx_msg, n_files = await asyncio.to_thread(_do_delete_tx)
        if tx_err != ErrorCode.SUCCESS:
            self.log.error(f"delete_kb tx: {tx_msg}")
            return tx_err, tx_msg

        self.log.info(
            f"Deleted knowledge base and cascade: id={kb_id}, "
            f"MySQL files={n_files} kb=1, Qdrant points={qdrant_deleted}, MinIO objects={minio_deleted}"
        )
        return ErrorCode.SUCCESS, None

    async def export_knowledge_base(
        self, req: KnowledgeBaseIdRequest, user_id: int
    ) -> Tuple[ErrorCode, object]:
        """
        导出知识库：返回知识库基础信息 + 关联文件列表（用于前端下载备份）。
        暂不导出向量数据，仅导出 MySQL 中的 rag_tbl / rag_file_tbl 元数据。
        """
        # 校验知识库归属
        err, row = await asyncio.to_thread(
            self.repo.get_knowledge_base_by_id, req.id, binding_user_id=user_id
        )
        if err != ErrorCode.SUCCESS or not row:
            msg = f"知识库不存在或无权导出: {req.id}"
            self.log.warning(msg)
            return ErrorCode.DATA_NOT_FOUND, msg

        kb = self._row_to_kb(row)

        # 读取部分文件记录用于导出（最多 1000 条，后续可根据需要扩展为分页导出）
        err, file_rows = await asyncio.to_thread(
            self.repo.list_files_by_kb, req.id, limit=1000
        )
        if err != ErrorCode.SUCCESS:
            self.log.error(f"list_files_by_kb failed when exporting kb: {err}")
            return err, "查询知识库文件列表失败"

        files: List[Dict[str, Any]] = []
        for r in file_rows:
            upload_time = r.get("upload_time")
            upload_ts = int(upload_time.timestamp() * 1000) if upload_time is not None else 0
            files.append(
                {
                    "id": r.get("id"),
                    "file_name": r.get("file_name", ""),
                    "file_path": r.get("file_path", ""),
                    "extension": r.get("extension", ""),
                    "status": int(r.get("status", 1) or 1),
                    "upload_time": upload_ts,
                    "file_char_count": int(r.get("file_char_count", 0) or 0),
                    "recall_count": int(r.get("recall_count", 0) or 0),
                    "rag_id": int(r.get("rag_id", 0) or 0),
                    "segmentation_type": r.get("segmentation_type", ""),
                }
            )

        return ErrorCode.SUCCESS, {"kb": kb.model_dump(), "files": files}

    # -------------------------------------------------------------------------
    # Neo4j Ontology Methods
    # -------------------------------------------------------------------------

    async def get_ontology_schema(self) -> Tuple[ErrorCode, object]:
        """
        获取本体Schema（标签和关系类型）
        """
        if not self._neo4j_client:
            return ErrorCode.EXTERNAL_SERVICE_UNAVAILABLE, "Neo4j service not available"

        # 获取所有标签
        err, labels_res = await asyncio.to_thread(
            self._neo4j_client.execute_query, "CALL db.labels()"
        )
        if err != ErrorCode.SUCCESS:
            return err, "Failed to fetch labels"
        labels = [row["label"] for row in labels_res]

        # 获取所有关系类型
        err, rels_res = await asyncio.to_thread(
            self._neo4j_client.execute_query, "CALL db.relationshipTypes()"
        )
        if err != ErrorCode.SUCCESS:
            return err, "Failed to fetch relationship types"
        relationship_types = [row["relationshipType"] for row in rels_res]

        return ErrorCode.SUCCESS, {
            "labels": labels,
            "relationship_types": relationship_types
        }


    async def get_ontology_nodes(self, label: str = None, page: int = 1, page_size: int = 50, keyword: str = None) -> Tuple[ErrorCode, object]:
        """
        获取本体节点列表，支持按标签筛选和关键词搜索
        """
        if not self._neo4j_client:
            return ErrorCode.EXTERNAL_SERVICE_UNAVAILABLE, "Neo4j service not available"

        skip = (page - 1) * page_size
        params = {"skip": skip, "limit": page_size}
        
        # 构建Cypher查询
        where_clauses = []
        if label and label != "all":
            label_part = f":`{label}`"
        else:
            label_part = ""
            
        if keyword:
            # 简单模糊搜索 name 或 id 属性
            where_clauses.append("(n.name CONTAINS $keyword OR n.id CONTAINS $keyword OR n.code CONTAINS $keyword)")
            params["keyword"] = keyword

        where_str = " WHERE " + " AND ".join(where_clauses) if where_clauses else ""
        
        # 查询总数
        count_query = f"MATCH (n{label_part}){where_str} RETURN count(n) as total"
        err, count_res = await asyncio.to_thread(
            self._neo4j_client.execute_query, count_query, params
        )
        if err != ErrorCode.SUCCESS:
            return err, "Failed to count nodes"
        total = count_res[0]["total"] if count_res else 0

        # 查询数据
        # 使用 id(n) 获取内部ID
        data_query = f"MATCH (n{label_part}){where_str} RETURN n, id(n) as internal_id, labels(n) as labels SKIP $skip LIMIT $limit"
        err, data_res = await asyncio.to_thread(
            self._neo4j_client.execute_query, data_query, params
        )
        if err != ErrorCode.SUCCESS:
            return err, "Failed to fetch nodes"
            
        # 格式化结果
        items = []
        for row in data_res:
            node = row["n"]
            labels = row["labels"]
            internal_id = row["internal_id"]
            
            # 提取主要属性
            # 优先使用属性中的 id, code, name, title
            display_id = node.get("id") or node.get("code") or str(internal_id)
            display_name = node.get("name") or node.get("title") or f"Node {internal_id}"
            
            item = {
                "id": display_id,
                "internal_id": internal_id,
                "name": display_name,
                "labels": labels,
                "properties": node,
                "created_at": node.get("created_at") or node.get("create_time"),
                "updated_at": node.get("updated_at") or node.get("update_time")
            }
            # Determine a primary label for display
            primary_label = labels[0] if labels else "Unknown"
            item["type"] = primary_label
            items.append(item)

        return ErrorCode.SUCCESS, {
            "total": total,
            "page": page,
            "page_size": page_size,
            "items": items
        }

    async def get_ontology_statistics(self) -> Tuple[ErrorCode, object]:
        """
        获取本体统计信息（节点总数、关系总数、类型数）
        返回格式：{ "total_nodes": 节点总数, "total_relations": 关系总数, "labels": {label: count} }
        """
        if not self._neo4j_client:
            self.log.error("get_ontology_statistics: Neo4j client is None")
            # Return empty stats instead of error to avoid frontend breaking
            return ErrorCode.SUCCESS, {"total_nodes": 0, "total_relations": 0, "labels": {}}
        
        try:
            # 查询节点总数
            node_query = "MATCH (n) RETURN count(n) as node_count"
            err, node_res = await asyncio.to_thread(
                self._neo4j_client.execute_query, node_query
            )
            if err != ErrorCode.SUCCESS:
                self.log.error(f"get_ontology_statistics: query nodes failed err={err}")
                node_count = 0
            else:
                node_count = node_res[0]["node_count"] if node_res else 0
            
            # 查询关系总数
            rel_query = "MATCH ()-[r]->() RETURN count(r) as rel_count"
            err, rel_res = await asyncio.to_thread(
                self._neo4j_client.execute_query, rel_query
            )
            if err != ErrorCode.SUCCESS:
                self.log.error(f"get_ontology_statistics: query rels failed err={err}")
                rel_count = 0
            else:
                rel_count = rel_res[0]["rel_count"] if rel_res else 0
            
            # 查询所有标签（用于计算类型数）
            label_query = "CALL db.labels() YIELD label RETURN label"
            err, label_res = await asyncio.to_thread(
                self._neo4j_client.execute_query, label_query
            )
            
            labels_map = {}
            if err != ErrorCode.SUCCESS:
                self.log.error(f"get_ontology_statistics: query labels failed err={err}")
            else:
                # 构造 labels 字典，目前仅作为类型计数，暂不统计每个 label 的节点数（避免性能开销）
                if label_res:
                    for record in label_res:
                        lbl = record.get("label")
                        if lbl:
                            labels_map[lbl] = 0
            
            # 如果db.labels()失败或为空，尝试另一种方式
            if not labels_map and err != ErrorCode.SUCCESS:
                 type_query = "MATCH (n) RETURN count(DISTINCT labels(n)) as type_count"
                 err, type_res = await asyncio.to_thread(
                    self._neo4j_client.execute_query, type_query
                )
                 # 如果只能拿到数量，就构造假label或者只返回数量字段（取决于前端需求，前端目前只用keys长度）
                 # 这里保持空字典，前端显示0
                 pass
            
            # 查询总召回次数
            err_recall, total_recall = await asyncio.to_thread(self.repo.get_total_recall_count)
            if err_recall != ErrorCode.SUCCESS:
                self.log.error(f"get_ontology_statistics: get_total_recall_count failed err={err_recall}")
                total_recall = 0

            statistics = {
                "total_nodes": node_count,
                "total_relations": rel_count,
                "labels": labels_map,
                "total_recall_count": total_recall,
                # 兼容旧字段
                "nodes": node_count,
                "relationships": rel_count,
                "types": len(labels_map)
            }
            
            self.log.info(f"Ontology statistics: nodes={node_count}, relationships={rel_count}, types={len(labels_map)}, recall={total_recall}")
            return ErrorCode.SUCCESS, statistics
            
        except Exception as e:
            self.log.error(f"get_ontology_statistics exception: {e}", exc_info=True)
            return ErrorCode.SUCCESS, {"total_nodes": 0, "total_relations": 0, "labels": {}}

    async def create_ontology_node(self, req: OntologyNodeCreateRequest) -> Tuple[ErrorCode, object]:
        """
        创建本体节点
        """
        if not self._neo4j_client:
            return ErrorCode.EXTERNAL_SERVICE_UNAVAILABLE, "Neo4j service not available"

        query = (
            f"CREATE (n:`{req.label}` $props) "
            "RETURN n, id(n) as internal_id"
        )
        params = {"props": req.properties}
        # 如果 name 不在 properties 中，手动加上
        if "name" not in params["props"]:
            params["props"]["name"] = req.name
            
        err, res = await asyncio.to_thread(
            self._neo4j_client.execute_query, query, params
        )
        if err != ErrorCode.SUCCESS:
            return err, "Failed to create node"
            
        if not res:
            return ErrorCode.OPERATION_FAILED, "Node created but no result returned"
            
        node = res[0]["n"]
        internal_id = res[0]["internal_id"]
        
        return ErrorCode.SUCCESS, {
            "id": str(internal_id),
            "internal_id": internal_id,
            "labels": [req.label],
            "properties": node
        }

    async def update_ontology_node(self, req: OntologyNodeUpdateRequest) -> Tuple[ErrorCode, object]:
        """
        更新本体节点
        """
        if not self._neo4j_client:
            return ErrorCode.EXTERNAL_SERVICE_UNAVAILABLE, "Neo4j service not available"

        # 根据 internal_id 查找
        try:
            node_id = int(req.id)
        except ValueError:
            return ErrorCode.INVALID_PARAMETER, "Invalid node ID format"

        set_clauses = []
        params = {"id": node_id}
        
        if req.name:
            set_clauses.append("n.name = $name")
            params["name"] = req.name
            
        if req.properties:
            for k, v in req.properties.items():
                set_clauses.append(f"n.`{k}` = ${k}")
                params[k] = v
        
        if not set_clauses:
            return ErrorCode.SUCCESS, None

        set_str = " SET " + ", ".join(set_clauses)
        query = f"MATCH (n) WHERE id(n) = $id {set_str} RETURN n"
        
        err, res = await asyncio.to_thread(
            self._neo4j_client.execute_query, query, params
        )
        if err != ErrorCode.SUCCESS:
            return err, "Failed to update node"
            
        return ErrorCode.SUCCESS, None

    async def delete_ontology_node(self, req: OntologyNodeDeleteRequest) -> Tuple[ErrorCode, object]:
        """
        删除本体节点
        """
        if not self._neo4j_client:
            return ErrorCode.EXTERNAL_SERVICE_UNAVAILABLE, "Neo4j service not available"

        try:
            node_id = int(req.id)
        except ValueError:
            return ErrorCode.INVALID_PARAMETER, "Invalid node ID format"
            
        # DETACH DELETE 删除节点及其关系
        query = "MATCH (n) WHERE id(n) = $id DETACH DELETE n"
        params = {"id": node_id}
        
        err, _ = await asyncio.to_thread(
            self._neo4j_client.execute_query, query, params
        )
        if err != ErrorCode.SUCCESS:
            return err, "Failed to delete node"
            
        return ErrorCode.SUCCESS, None

    async def update_file_status(
        self, req: FileUpdateStatusRequest, user_id: int
    ) -> Tuple[ErrorCode, object]:
        """更新单个文件的启用状态，并同步刷新知识库容量与文件数。"""
        # 校验知识库归属
        err, kb_row = await asyncio.to_thread(
            self.repo.get_knowledge_base_by_id, req.kb_id, binding_user_id=user_id
        )
        if err != ErrorCode.SUCCESS or not kb_row:
            msg = f"知识库不存在或无权操作: {req.kb_id}"
            self.log.warning(msg)
            return ErrorCode.DATA_NOT_FOUND, msg

        # 查询文件记录
        err, file_row = await asyncio.to_thread(self.repo.get_file_by_id, req.file_id)
        if err != ErrorCode.SUCCESS or not file_row:
            msg = f"文件不存在: {req.file_id}"
            self.log.warning(msg)
            return ErrorCode.DATA_NOT_FOUND, msg
        if int(file_row.get("rag_id", 0) or 0) != req.kb_id:
            msg = f"文件不属于指定知识库: file_id={req.file_id}, kb_id={req.kb_id}"
            self.log.warning(msg)
            return ErrorCode.DATA_INVALID, msg

        new_status = 1 if req.status == 1 else 0
        old_status = int(file_row.get("status", 1) or 1)
        if new_status == old_status:
            return ErrorCode.SUCCESS, None

        # 更新文件状态
        upd_err, _ = await asyncio.to_thread(
            self.repo.update_file_record, req.file_id, {"status": new_status}
        )
        if upd_err != ErrorCode.SUCCESS:
            self.log.error(f"update_file_status: update_file_record failed err={upd_err}, file_id={req.file_id}")
            return upd_err, "更新文件状态失败"

        # 若文件已完成嵌入，则重新统计知识库的容量与文件数
        embedding_process = int(file_row.get("embedding_process", 0) or 0)
        if embedding_process == 3:
            try:
                stats_err, (total_chars, file_cnt) = await asyncio.to_thread(
                    self.repo.calc_kb_file_stats, req.kb_id
                )
                if stats_err == ErrorCode.SUCCESS:
                    upd_kb_err, _ = await asyncio.to_thread(
                        self.repo.update_kb_file_stats, req.kb_id, total_chars, file_cnt
                    )
                    if upd_kb_err != ErrorCode.SUCCESS:
                        self.log.warning(
                            f"update_file_status: update_kb_file_stats failed: {upd_kb_err} for kb_id={req.kb_id}"
                        )
                else:
                    self.log.warning(
                        f"update_file_status: calc_kb_file_stats failed: {stats_err} for kb_id={req.kb_id}"
                    )
            except Exception as e:
                self.log.warning(f"update_file_status: refresh kb stats exception for kb_id={req.kb_id}: {e}")

        return ErrorCode.SUCCESS, None

    async def delete_file(
        self, req: FileIdWithKbRequest, user_id: int
    ) -> Tuple[ErrorCode, object]:
        """删除单个文件：清理 Redis/Qdrant/MinIO/MySQL，并刷新知识库容量与文件数。"""
        # 校验知识库归属
        err, kb_row = await asyncio.to_thread(
            self.repo.get_knowledge_base_by_id, req.kb_id, binding_user_id=user_id
        )
        if err != ErrorCode.SUCCESS or not kb_row:
            msg = f"知识库不存在或无权操作: {req.kb_id}"
            self.log.warning(msg)
            return ErrorCode.DATA_NOT_FOUND, msg

        # 查询文件记录
        err, file_row = await asyncio.to_thread(self.repo.get_file_by_id, req.file_id)
        if err != ErrorCode.SUCCESS or not file_row:
            msg = f"文件不存在: {req.file_id}"
            self.log.warning(msg)
            return ErrorCode.DATA_NOT_FOUND, msg
        if int(file_row.get("rag_id", 0) or 0) != req.kb_id:
            msg = f"文件不属于指定知识库: file_id={req.file_id}, kb_id={req.kb_id}"
            self.log.warning(msg)
            return ErrorCode.DATA_INVALID, msg

        object_name = file_row.get("file_path") or None
        # 清理单个文件相关资源（内部会删除 rag_file_tbl 记录）
        await asyncio.to_thread(
            self.cleanup_file_resources, req.kb_id, req.file_id, object_name, None
        )

        # 删除后重新统计知识库的容量与文件数
        try:
            stats_err, (total_chars, file_cnt) = await asyncio.to_thread(
                self.repo.calc_kb_file_stats, req.kb_id
            )
            if stats_err == ErrorCode.SUCCESS:
                upd_kb_err, _ = await asyncio.to_thread(
                    self.repo.update_kb_file_stats, req.kb_id, total_chars, file_cnt
                )
                if upd_kb_err != ErrorCode.SUCCESS:
                    self.log.warning(
                        f"delete_file: update_kb_file_stats failed: {upd_kb_err} for kb_id={req.kb_id}"
                    )
            else:
                self.log.warning(
                    f"delete_file: calc_kb_file_stats failed: {stats_err} for kb_id={req.kb_id}"
                )
        except Exception as e:
            self.log.warning(f"delete_file: refresh kb stats exception for kb_id={req.kb_id}: {e}")

        return ErrorCode.SUCCESS, None

    async def get_file_download_url(
        self, req: FileIdWithKbRequest, user_id: int
    ) -> Tuple[ErrorCode, object]:
        """获取文件下载链接（MinIO 预签名 URL）。"""
        if not self._minio_client:
            self.log.warning("get_file_download_url: MinIO client not configured")
            return ErrorCode.SERVICE_UNAVAILABLE, "文件存储服务未配置"

        # 校验知识库归属
        err, kb_row = await asyncio.to_thread(
            self.repo.get_knowledge_base_by_id, req.kb_id, binding_user_id=user_id
        )
        if err != ErrorCode.SUCCESS or not kb_row:
            msg = f"知识库不存在或无权操作: {req.kb_id}"
            self.log.warning(msg)
            return ErrorCode.DATA_NOT_FOUND, msg

        # 查询文件记录
        err, file_row = await asyncio.to_thread(self.repo.get_file_by_id, req.file_id)
        if err != ErrorCode.SUCCESS or not file_row:
            msg = f"文件不存在: {req.file_id}"
            self.log.warning(msg)
            return ErrorCode.DATA_NOT_FOUND, msg
        if int(file_row.get("rag_id", 0) or 0) != req.kb_id:
            msg = f"文件不属于指定知识库: file_id={req.file_id}, kb_id={req.kb_id}"
            self.log.warning(msg)
            return ErrorCode.DATA_INVALID, msg

        object_name = file_row.get("file_path")
        if not object_name:
            self.log.warning(f"get_file_download_url: file_path empty for file_id={req.file_id}")
            return ErrorCode.DATA_INVALID, "文件存储路径不存在"

        try:
            expires = timedelta(hours=1)
            file_name = file_row.get("file_name") or object_name.split("/")[-1]
            encoded_name = quote(file_name)
            response_headers = {
                "response-content-disposition": f"attachment; filename*=UTF-8''{encoded_name}"
            }
            url = await asyncio.to_thread(
                self._minio_client.presigned_get_object,
                self._minio_bucket,
                object_name,
                expires,
                response_headers,
            )
        except Exception as e:
            self.log.error(f"get_file_download_url: presigned_get_object failed for file_id={req.file_id}: {e}")
            return ErrorCode.EXTERNAL_SERVICE_ERROR, "生成下载地址失败"

        data = {
            "url": url,
            "file_name": file_row.get("file_name", ""),
        }
        return ErrorCode.SUCCESS, data

    async def get_knowledge_base(
        self, req: KnowledgeBaseIdRequest, user_id: int
    ) -> Tuple[ErrorCode, object]:
        """查询单个知识库详情（仅当前用户绑定的知识库）"""
        err, row = await asyncio.to_thread(
            self.repo.get_knowledge_base_by_id, req.id, binding_user_id=user_id
        )
        if err != ErrorCode.SUCCESS or not row:
            msg = f"知识库不存在或无权查看: {req.id}"
            self.log.warning(msg)
            return ErrorCode.DATA_NOT_FOUND, msg

        kb = self._row_to_kb(row)
        return ErrorCode.SUCCESS, kb

    async def list_knowledge_items(
        self, req: KnowledgeItemListRequest, user_id: int
    ) -> Tuple[ErrorCode, object]:
        err, kb_row = await asyncio.to_thread(
            self.repo.get_knowledge_base_by_id, req.kb_id, binding_user_id=user_id
        )
        if err != ErrorCode.SUCCESS or not kb_row:
            msg = f"知识库不存在或无权查看: {req.kb_id}"
            self.log.warning(msg)
            return ErrorCode.DATA_NOT_FOUND, msg

        err, rows = await asyncio.to_thread(
            self.repo.list_knowledge_items_by_kb,
            req.kb_id,
            req.min_recall_count,
            req.max_recall_count,
            req.recent_days,
            req.sort_by,
            req.sort_order or "desc",
        )
        if err != ErrorCode.SUCCESS:
            self.log.error(f"list_knowledge_items_by_kb failed: {err} for kb_id={req.kb_id}")
            return err, "查询知识信息列表失败"

        items: List[KnowledgeItemInfo] = []
        for r in rows:
            upload_time = r.get("upload_time")
            ts = (
                upload_time.strftime("%Y-%m-%d %H:%M:%S")
                if upload_time is not None
                else None
            )
            item = KnowledgeItemInfo(
                id=int(r.get("id", 0) or 0),
                file_name=r.get("file_name", "") or "",
                recall_count=int(r.get("recall_count", 0) or 0),
                last_access_time=None,
                avg_score=None,
                instruction_score=None,
                relevance_score=None,
                is_noise=False,
                is_redundant=False,
                created_at=ts,
                updated_at=ts,
            )
            items.append(item)

        return ErrorCode.SUCCESS, {"items": [i.model_dump() for i in items]}

    async def list_knowledge_bases(
        self, req: KnowledgeBaseListRequest, user_id: int
    ) -> Tuple[ErrorCode, object]:
        """分页查询当前用户知识库列表 -> rag_tbl"""
        err, total = await asyncio.to_thread(
            self.repo.count_knowledge_bases, binding_user_id=user_id, keyword=req.keyword
        )
        if err != ErrorCode.SUCCESS:
            self.log.error(f"count rag_tbl failed: {err}")
            return err, "查询知识库数量失败"

        offset = (req.page_no - 1) * req.page_size
        err, rows = await asyncio.to_thread(
            self.repo.list_knowledge_bases,
            binding_user_id=user_id,
            keyword=req.keyword,
            limit=req.page_size,
            offset=offset,
        )
        if err != ErrorCode.SUCCESS:
            self.log.error(f"select rag_tbl failed: {err}")
            return err, "查询知识库列表失败"

        items: List[KnowledgeBaseInfo] = [self._row_to_kb(r) for r in rows]
        page = KnowledgeBasePage(
            items=items,
            total=total,
            page_no=req.page_no,
            page_size=req.page_size,
        )
        return ErrorCode.SUCCESS, page

    # ============== 知识片段 ==============
    async def list_segments(
        self, req: KnowledgeSegmentListRequest, user_id: int
    ) -> Tuple[ErrorCode, object]:
        """
        查询指定文件下的知识片段列表。
        """
        err_file, file_row = await asyncio.to_thread(self.repo.get_file_by_id, req.file_id)
        if err_file != ErrorCode.SUCCESS or not file_row:
            return ErrorCode.DATA_NOT_FOUND, "文件不存在"
        kb_id = int(file_row.get("rag_id", 0) or 0)
        err_kb, kb_row = await asyncio.to_thread(
            self.repo.get_knowledge_base_by_id, kb_id, binding_user_id=user_id
        )
        if err_kb != ErrorCode.SUCCESS or not kb_row:
            return ErrorCode.PERMISSION_DENIED, "无权操作该知识库"
        limit = req.page_size
        offset = (req.page_no - 1) * req.page_size
        err, rows = await asyncio.to_thread(
            self.repo.list_segments_by_file, 
            req.file_id, 
            req.status, 
            user_id, 
            limit, 
            offset,
            req.sort_by,
            req.sort_order
        )
        if err != ErrorCode.SUCCESS:
            return err, "查询知识片段失败"
        items: List[KnowledgeSegmentInfo] = []
        for r in rows:
            items.append(
                KnowledgeSegmentInfo(
                    id=int(r.get("id", 0) or 0),
                    title=r.get("title", "") or "",
                    content=r.get("content"),
                    summary=r.get("summary"),
                    snippet_type=r.get("snippet_type") or "text",
                    sort_index=int(r.get("sort_index", 9999) or 9999),
                    status=int(r.get("status", 1) or 1),
                    recall_count=int(r.get("recall_count", 0) or 0),
                    last_search_time=(
                        r.get("last_search_time").strftime("%Y-%m-%d %H:%M:%S")
                        if r.get("last_search_time") is not None
                        else None
                    ),
                    average_retrieval_score=float(r.get("average_retrieval_score", 0.0) or 0.0) if r.get("average_retrieval_score") is not None else None,
                    command_score=float(r.get("command_score", 0.0) or 0.0) if r.get("command_score") is not None else None,
                    relevance_score=float(r.get("relevance_score", 0.0) or 0.0) if r.get("relevance_score") is not None else None,
                    is_noise=bool(r.get("is_noise", False)),
                    is_redundant=bool(r.get("is_redundant", False)),
                    created_at=(
                        r.get("create_time").strftime("%Y-%m-%d %H:%M:%S")
                        if r.get("create_time") is not None
                        else None
                    ),
                    updated_at=(
                        r.get("update_time").strftime("%Y-%m-%d %H:%M:%S")
                        if r.get("update_time") is not None
                        else None
                    ),
                )
            )
        return ErrorCode.SUCCESS, {"items": [i.model_dump() for i in items]}

    async def create_knowledge_segment(
        self, req: KnowledgeSegmentCreateRequest, user_id: int
    ) -> Tuple[ErrorCode, object]:
        """创建知识片段"""
        # 1. 获取文件信息
        err, file_row = await asyncio.to_thread(self.repo.get_file_by_id, req.file_id)
        if err != ErrorCode.SUCCESS or not file_row:
            return ErrorCode.DATA_NOT_FOUND, "文件不存在"
        
        kb_id = int(file_row.get("rag_id", 0))
        # 校验权限
        err_kb, kb_row = await asyncio.to_thread(self.repo.get_knowledge_base_by_id, kb_id, binding_user_id=user_id)
        if err_kb != ErrorCode.SUCCESS or not kb_row:
             return ErrorCode.PERMISSION_DENIED, "无权操作该知识库"

        # 2. 获取最大 sort_index
        err, max_idx = await asyncio.to_thread(self.repo.get_max_sort_index, req.file_id)
        if err != ErrorCode.SUCCESS:
            return err, "获取排序索引失败"
        new_sort_index = max_idx + 1

        # 3. 插入 MySQL
        point_id = None
        if req.status == 2:
            point_id = (kb_id << 40) | (req.file_id << 20) | new_sort_index

        seg_data = {
            "file_id": req.file_id,
            "binding_user_id": user_id,
            "title": req.title or f"Segment#{new_sort_index + 1}",
            "content": req.content,
            "summary": req.content[:200],
            "sort_index": new_sort_index,
            "status": req.status,
            "snippet_type": "text",
            "recall_count": 0,
            "qdrant_vector_id": str(point_id) if point_id is not None else None,
        }
        err, seg_id = await asyncio.to_thread(self.repo.insert_segment, seg_data)
        if err != ErrorCode.SUCCESS:
            return err, "创建知识片段失败"

        # 4. 如果状态是启用，插入 Qdrant
        if req.status == 2 and self._qdrant and self._embedding_client:
            try:
                # 计算向量
                err_embed, embeddings = await asyncio.to_thread(self.embed_texts, [req.content])
                if err_embed == ErrorCode.SUCCESS and embeddings:
                    vec = embeddings[0]
                    # point_id 已在上方计算
                    point = {
                        "id": point_id,
                        "vector": vec,
                        "payload": {
                            "rag_id": kb_id,
                            "file_id": req.file_id,
                            "segment_id": seg_id,
                            "chunk_index": new_sort_index,
                            "text": req.content[:2000],
                            "binding_user_id": user_id,
                        },
                    }
                    await asyncio.to_thread(self._qdrant.insert, [point])
                    self.log.info(f"create_knowledge_segment: inserted Qdrant point {point_id}")
            except Exception as e:
                self.log.warning(f"create_knowledge_segment: Qdrant insert failed: {e}")

        return ErrorCode.SUCCESS, {"id": seg_id}

    async def update_knowledge_segment(
        self, req: KnowledgeSegmentUpdateRequest, user_id: int
    ) -> Tuple[ErrorCode, object]:
        """更新知识片段信息（内容、状态等），并同步更新 Qdrant 向量/Payload"""
        # 1. 检查片段是否存在
        err, segment = await asyncio.to_thread(self.repo.get_segment_by_id, req.id)
        if err != ErrorCode.SUCCESS or not segment:
            return ErrorCode.DATA_NOT_FOUND, "知识片段不存在"
            
        # 2. 准备更新数据
        data = {}
        old_status = int(segment.get("status", 1))
        old_content = segment.get("content", "")
        
        new_content = req.content if req.content is not None else old_content
        new_status = req.status if req.status is not None else old_status
        
        content_changed = (req.content is not None and req.content != old_content)
        status_changed = (req.status is not None and req.status != old_status)

        if content_changed:
            data["content"] = req.content
        if req.status is not None:
            data["status"] = req.status
        if req.title is not None:
            data["title"] = req.title
            
        if not data:
            return ErrorCode.SUCCESS, None

        # 3. 同步 Qdrant 向量
        # 获取关联信息
        file_id = segment.get("file_id")
        sort_index = segment.get("sort_index")
        
        # 获取文件信息 (rag_id)
        err_file, file_row = await asyncio.to_thread(self.repo.get_file_by_id, file_id)
        kb_id = int(file_row.get("rag_id", 0)) if file_row else 0
        err_kb, kb_row = await asyncio.to_thread(
            self.repo.get_knowledge_base_by_id, kb_id, binding_user_id=user_id
        )
        if err_kb != ErrorCode.SUCCESS or not kb_row:
            return ErrorCode.PERMISSION_DENIED, "无权操作该知识库"
        
        if self._qdrant and kb_id > 0:
            # Case A: 需要在 Qdrant 中存在（启用状态）且（内容变更 OR 状态刚变为启用）
            # 注意：如果状态一直是启用，但内容变更了，也需要更新
            should_upsert = (new_status == 2) and (content_changed or (status_changed and old_status != 2))
            
            # Case B: 需要从 Qdrant 中删除（状态变为非启用，且之前可能是启用或存在脏数据）
            should_delete = (new_status != 2) and (old_status == 2)
            
            if should_upsert and self._embedding_client:
                # 计算新向量
                err_embed, embeddings = await asyncio.to_thread(self.embed_texts, [new_content])
                if err_embed == ErrorCode.SUCCESS and embeddings:
                    vec = embeddings[0]
                    point_id = (kb_id << 40) | (file_id << 20) | sort_index
                    point = {
                        "id": point_id,
                        "vector": vec,
                        "payload": {
                            "rag_id": kb_id,
                            "file_id": file_id,
                            "segment_id": req.id,
                            "chunk_index": sort_index,
                            "text": new_content[:2000],
                            "binding_user_id": user_id,
                        },
                    }
                    asyncio.create_task(asyncio.to_thread(self._qdrant.insert, [point]))
                    self.log.info(f"update_knowledge_segment: upsert Qdrant point {point_id} for segment {req.id}")

            elif should_delete:
                point_id = (kb_id << 40) | (file_id << 20) | sort_index
                asyncio.create_task(asyncio.to_thread(self._qdrant.delete, ids=[point_id], collection_name=self._rag_collection))
                self.log.info(f"update_knowledge_segment: deleted Qdrant point {point_id} due to status change")
            
        # 4. 执行 MySQL 更新
        err, _ = await asyncio.to_thread(self.repo.update_segment, req.id, data)
        if err != ErrorCode.SUCCESS:
            return err, "更新知识片段失败"
        
        return ErrorCode.SUCCESS, None

    async def delete_knowledge_segment(
        self, req: KnowledgeSegmentIdRequest, user_id: int
    ) -> Tuple[ErrorCode, object]:
        """删除知识片段，并同步删除 Qdrant 向量"""
        # 1. 检查片段是否存在
        err, segment = await asyncio.to_thread(self.repo.get_segment_by_id, req.id)
        if err != ErrorCode.SUCCESS or not segment:
            return ErrorCode.DATA_NOT_FOUND, "知识片段不存在"
            
        file_id = segment.get("file_id")
        sort_index = segment.get("sort_index")
        
        # 获取文件信息 (rag_id)
        err_file, file_row = await asyncio.to_thread(self.repo.get_file_by_id, file_id)
        if err_file == ErrorCode.SUCCESS and file_row:
            kb_id = int(file_row.get("rag_id", 0))
            err_kb, kb_row = await asyncio.to_thread(
                self.repo.get_knowledge_base_by_id, kb_id, binding_user_id=user_id
            )
            if err_kb != ErrorCode.SUCCESS or not kb_row:
                return ErrorCode.PERMISSION_DENIED, "无权操作该知识库"
            
            # 删除 Qdrant 向量
            if self._qdrant:
                # 构造 point_id
                point_id = (kb_id << 40) | (file_id << 20) | sort_index
                try:
                    await asyncio.to_thread(self._qdrant.delete, ids=[point_id], collection_name=self._rag_collection)
                    self.log.info(f"delete_knowledge_segment: deleted Qdrant point {point_id} for segment {req.id}")
                except Exception as e:
                    self.log.warning(f"delete_knowledge_segment: Qdrant delete failed: {e}")
        
        # 2. 删除 MySQL 记录
        err, _ = await asyncio.to_thread(self.repo.delete_segment, req.id)
        if err != ErrorCode.SUCCESS:
            return err, "删除知识片段失败"
            
        return ErrorCode.SUCCESS, None

    # ============== RAG 问答（使用 vllm Embedding + Qdrant） ==============

    async def rag_query(
        self, req: RagQueryRequest, user_id: int
    ) -> Tuple[ErrorCode, object]:
        """RAG 问答接口：对 query 做向量化，在 Qdrant 中按知识库与用户过滤做相似检索。"""
        # 1. 校验知识库归属
        err, kb_row = await asyncio.to_thread(
            self.repo.get_knowledge_base_by_id,
            req.kb_id, binding_user_id=user_id
        )
        if err != ErrorCode.SUCCESS or not kb_row:
            msg = f"知识库不存在: {req.kb_id}"
            self.log.warning(msg)
            return ErrorCode.DATA_NOT_FOUND, msg

        kb = self._row_to_kb(kb_row)

        # 2. 检查向量服务是否可用
        if not self._embedding_client or not self._embedding_client.is_available() or not self._qdrant:
            msg = "RAG 向量服务未启用或不可用"
            self.log.warning(f"rag_query: {msg}")
            return ErrorCode.SERVICE_UNAVAILABLE, msg

        # 3. 对查询语句做向量化（embed 为同步 I/O，放入线程）
        err, embeddings = await asyncio.to_thread(self.embed_texts, [req.query])
        if err != ErrorCode.SUCCESS or not embeddings or not embeddings[0]:
            self.log.error(f"rag_query: embed_texts failed err={err}")
            return err if err != ErrorCode.SUCCESS else ErrorCode.DATA_INVALID, "问题向量化失败"
        query_vec = embeddings[0]

        # 4. 在 Qdrant 中按 rag_id / binding_user_id 过滤做相似度检索
        qdrant_filter = {
            "must": [
                {"key": "rag_id", "match": {"value": kb.id}},
                {"key": "binding_user_id", "match": {"value": user_id}},
            ]
        }
        err, results = await asyncio.to_thread(
            self._qdrant.query_by_vector_with_filter,
            query_vector=query_vec,
            collection_name=self._rag_collection,
            qdrant_filter=qdrant_filter,
            limit=req.top_k,
            score_threshold=0.5,
            with_payload=True,
        )
        if err != ErrorCode.SUCCESS:
            self.log.error(f"rag_query: Qdrant query failed err={err}")
            return err, "知识检索失败"

        # 5. 构造来源文档（使用存储的分段文本片段）
        source_documents: List[str] = []
        
        # 从 Qdrant 结果中提取 segment_id
        segment_ids = []
        for item in results:
            payload = item.get("payload") or {}
            seg_id = payload.get("segment_id")
            if seg_id:
                try:
                    seg_id = int(seg_id)
                    segment_ids.append(seg_id)
                except (ValueError, TypeError):
                    pass
        
        # 批量查询 MySQL 获取最新内容
        mysql_segments = {}
        if segment_ids:
            err_seg, seg_rows = await asyncio.to_thread(self.repo.get_segments_by_ids, segment_ids)
            if err_seg == ErrorCode.SUCCESS:
                for row in seg_rows:
                    mysql_segments[row["id"]] = row
                
                # 异步增加召回次数
                asyncio.create_task(asyncio.to_thread(self.repo.increment_segment_recall, segment_ids))
            else:
                self.log.warning(f"rag_query: get_segments_by_ids failed: {err_seg}")

        # 组装结果
        for item in results:
            payload = item.get("payload") or {}
            seg_id = payload.get("segment_id")
            
            text = ""
            status = 2  # 默认启用（兼容旧数据）
            
            # 优先使用 MySQL 中的内容和状态
            if seg_id:
                try:
                    seg_id = int(seg_id)
                    if seg_id in mysql_segments:
                        seg_row = mysql_segments[seg_id]
                        text = seg_row.get("content", "")
                        status = int(seg_row.get("status", 2))
                    else:
                        # MySQL 中不存在（可能已物理删除），视为无效
                        continue
                except:
                    pass
            
            # 如果状态不是启用 (2)，则跳过
            if status != 2:
                continue

            # 降级使用 payload 中的 text
            if not text:
                text = payload.get("text", "")
            
            score = item.get("score", 0.0)
            snippet = text[:200].replace("\n", " ")
            source_documents.append(f"{snippet} ... (score={score:.3f})" if snippet else f"(score={score:.3f})")

        if not source_documents:
            answer = f"在知识库「{kb.name}」中未找到与问题「{req.query}」高度相关的内容。"
        else:
            # 简单拼接检索到的片段，后续可以接入大模型生成自然语言回答
            joined = "\n".join(f"- {doc}" for doc in source_documents)
            answer = (
                f"这是根据知识库「{kb.name}」中检索到的相关内容给出的回答示例：\n\n"
                f"{joined}\n\n"
                f"（后续可以在此基础上接入大模型生成更自然的回答）"
            )

        result = RagQueryAnswer(
            answer=answer,
            kb_id=kb.id,
            query=req.query,
            top_k=req.top_k,
            source_documents=source_documents,
        )
        return ErrorCode.SUCCESS, result

    async def all_kb_name_list(
        self, req: KnowledgeBaseListRequest, user_id: int
    ) -> Tuple[ErrorCode, object]:
        """查询当前用户下所有知识库名称列表（用于下拉等），返回 id 和 name 列表"""
        err, rows = await asyncio.to_thread(
            self.repo.list_knowledge_bases,
            binding_user_id=user_id,
            keyword=req.keyword,
            limit=1000,
            offset=0,
        )
        if err != ErrorCode.SUCCESS:
            return err, "查询知识库列表失败"
        items = [{"id": r.get("id"), "name": r.get("rag_name", "")} for r in rows if r.get("rag_name")]
        return ErrorCode.SUCCESS, items

    def get_file_vectorization_status(
        self, req: FileVectorizationStatusRequest, user_id: int
    ) -> Tuple[ErrorCode, FileVectorizationStatusResponse]:
        """
        查询文件向量化状态。支持按 file_id 或 file_name 查询。
        - 校验知识库归属
        - 检查MySQL中的文件记录
        - 检查Qdrant中的向量化数据
        """
        if req.file_id is None and not (req.file_name and req.file_name.strip()):
            return ErrorCode.INVALID_PARAMETER, FileVectorizationStatusResponse(
                file_name=req.file_name or "",
                kb_id=req.kb_id,
                is_vectorized=False,
                error_message="请提供 file_id 或 file_name"
            )

        # 1. 校验知识库归属
        err, kb_row = self.repo.get_knowledge_base_by_id(req.kb_id, binding_user_id=user_id)
        if err != ErrorCode.SUCCESS:
            msg = f"数据库查询失败: {err}"
            self.log.error(msg)
            return err, FileVectorizationStatusResponse(
                file_name=req.file_name or "",
                kb_id=req.kb_id,
                is_vectorized=False,
                error_message=msg
            )
        if not kb_row:
            msg = f"知识库不存在或无权操作: {req.kb_id}"
            self.log.warning(msg)
            return ErrorCode.DATA_NOT_FOUND, FileVectorizationStatusResponse(
                file_name=req.file_name or "",
                kb_id=req.kb_id,
                is_vectorized=False,
                error_message=msg
            )

        # 2. 获取文件记录：优先按 file_id，否则按 file_name + kb_id
        file_record = None
        if req.file_id is not None:
            err, file_record = self.repo.get_file_by_id(req.file_id)
            if err != ErrorCode.SUCCESS:
                self.log.error(f"get_file_vectorization_status: get_file_by_id failed err={err}")
                return err, FileVectorizationStatusResponse(
                    file_name=req.file_name or "",
                    kb_id=req.kb_id,
                    is_vectorized=False,
                    error_message="数据库查询失败"
                )
            if file_record and file_record.get("rag_id") != req.kb_id:
                return ErrorCode.DATA_NOT_FOUND, FileVectorizationStatusResponse(
                    file_name=file_record.get("file_name", ""),
                    kb_id=req.kb_id,
                    is_vectorized=False,
                    error_message="文件不属于该知识库"
                )
        else:
            err, file_record = self.repo.get_file_by_name_and_kb(req.file_name, req.kb_id)
            if err != ErrorCode.SUCCESS:
                self.log.error(f"get_file_vectorization_status: get_file_by_name_and_kb failed err={err}")
                return err, FileVectorizationStatusResponse(
                    file_name=req.file_name or "",
                    kb_id=req.kb_id,
                    is_vectorized=False,
                    error_message="数据库查询失败"
                )

        file_name_for_resp = (file_record.get("file_name") if file_record else None) or req.file_name or ""

        if not file_record:
            # MySQL中不存在文件记录
            return ErrorCode.SUCCESS, FileVectorizationStatusResponse(
                file_name=file_name_for_resp,
                kb_id=req.kb_id,
                is_vectorized=False,
                file_id=None,
                embedding_status=0,
                embedding_progress=0.0
            )

        # 3. 进度只查 Redis（每次写入 Qdrant 后已更新 Redis；完成时也在 Redis 标记）
        file_id = file_record.get("id")
        redis_progress = self.get_rag_file_progress(file_id)

        def _safe_float(v, default: float = 0.0) -> float:
            try:
                f = float(v) if v is not None else default
                return default if (f != f or f < 0 or f > 1) else min(f, 1.0)  # 过滤 NaN/Inf
            except (TypeError, ValueError):
                return default

        if redis_progress is not None:
            # 以 Redis 为准：进度、状态、是否已完成均来自 Redis
            embedding_process = redis_progress.get("embedding_process", 0)
            embedding_progress = _safe_float(redis_progress.get("embedding_progress"), 0.0)
            file_char_count = redis_progress.get("file_char_count") or 0
            file_embedding_offset = redis_progress.get("file_embedding_offset")
            if embedding_process != 3 and file_char_count and file_embedding_offset is not None:
                try:
                    embedding_progress = min(float(file_embedding_offset) / max(1, int(file_char_count)), 1.0)
                except (TypeError, ValueError, ZeroDivisionError):
                    pass
            error_message = redis_progress.get("error_message")
            is_vectorized = embedding_process == 3
            if is_vectorized:
                embedding_progress = 1.0
                self.clear_rag_file_progress(file_id)
        else:
            # 无 Redis 时用 DB 兜底（未启用 Redis 或尚未写入）
            embedding_process = file_record.get("embedding_process", 0)
            file_char_count = file_record.get("file_char_count", 0)
            file_embedding_offset = file_record.get("file_embedding_offset", 0)
            error_message = None
            if file_char_count and file_embedding_offset:
                try:
                    embedding_progress = min(float(file_embedding_offset) / max(1, int(file_char_count)), 1.0)
                except (TypeError, ValueError, ZeroDivisionError):
                    embedding_progress = 0.0
            else:
                embedding_progress = 0.0
            is_vectorized = embedding_process == 3
            if is_vectorized:
                embedding_progress = 1.0

        if embedding_process == 1:
            embedding_status = 2  # 正在向量化知识
        elif embedding_process == 2:
            embedding_status = 1  # 正在上传文件
        elif embedding_process == 4:
            embedding_status = 4  # 失败
        elif embedding_process == 3:
            embedding_status = 3  # 已完成
        else:
            embedding_status = 0  # 未开始

        return ErrorCode.SUCCESS, FileVectorizationStatusResponse(
            file_name=file_name_for_resp,
            kb_id=req.kb_id,
            is_vectorized=is_vectorized,
            file_id=file_id,
            embedding_status=embedding_status,
            embedding_progress=embedding_progress,
            error_message=error_message
        )

    # ============== 混合 RAG 查询（5阶段流程） ==============
    
    async def hybrid_rag_query(self, request: HybridRagQueryRequest) -> Tuple[ErrorCode, HybridRagQueryResponse]:
        """
        混合 RAG 查询：5阶段流程
        [阶段1: 向量初筛] → [阶段2: 元数据过滤] → [阶段3: 语义扩展] → [阶段4: 二次检索增强] → [阶段5: 结果映射与合并]
        """
        try:
            # 阶段1: 向量初筛 - 输入查询 → 向量化 → Qdrant相似性搜索 → 获取Top-N初筛结果
            initial_results = await self._stage1_vector_initial_screening(request)
            if not initial_results:
                return ErrorCode.SUCCESS, HybridRagQueryResponse(
                    original_query=request.query,
                    enhanced_query=request.query,
                    related_concepts=[],
                    results=[]
                )
            
            # 阶段2: 元数据过滤 - 初筛结果 → MySQL元数据/关键词过滤 → 过滤后结果
            filtered_results = await self._stage2_metadata_filtering(initial_results, request.kb_id)
            
            # 阶段3: 语义扩展 - 过滤结果 → 提取关键实体 → Neo4j图遍历/推理 → 获取相关概念
            related_concepts = await self._stage3_semantic_extension(request.query, filtered_results)
            
            # 阶段4: 二次检索增强 - 相关概念 + 原始查询 → 构造增强查询 → Qdrant二次检索 → 增强结果
            enhanced_results = await self._stage4_secondary_retrieval_enhancement(
                request.query, related_concepts, request.kb_id, request.top_k
            )
            
            # 合并初筛结果和增强结果
            combined_results = filtered_results + enhanced_results
            # 按分数降序排序，确保去重时保留高分结果
            combined_results.sort(key=lambda x: x.get("score", 0.0), reverse=True)
            
            # 阶段5: 结果映射与合并 - 增强结果 → 映射到MySQL知识片段 → 片段合并策略 → 最终输出
            final_results = await self._stage5_result_mapping_and_merging(combined_results, request.kb_id)
            
            # 构造增强查询文本
            enhanced_query = request.query
            if related_concepts:
                enhanced_query = f"{request.query} 相关概念：{', '.join(related_concepts[:5])}"
            
            return ErrorCode.SUCCESS, HybridRagQueryResponse(
                original_query=request.query,
                enhanced_query=enhanced_query,
                related_concepts=related_concepts,
                results=final_results
            )
            
        except Exception as e:
            self.log.error(f"hybrid_rag_query failed: {e}")
            return ErrorCode.SYSTEM_ERROR, HybridRagQueryResponse(
                original_query=request.query,
                enhanced_query=request.query,
                related_concepts=[],
                results=[]
            )
    
    async def _stage1_vector_initial_screening(self, request: HybridRagQueryRequest) -> List[Dict[str, Any]]:
        """阶段1: 向量初筛"""
        try:
            # 向量化查询
            err, embeddings = self.embed_texts([request.query])
            if err != ErrorCode.SUCCESS or not embeddings:
                self.log.warning(f"_stage1_vector_initial_screening: embed failed err={err}")
                return []
            
            query_vector = embeddings[0]
            
            # Qdrant 相似性搜索
            if not self._qdrant:
                self.log.warning("_stage1_vector_initial_screening: Qdrant not available")
                return []
            
            # 构建过滤器：仅搜索指定知识库
            filter_condition = rest.Filter(
                must=[
                    rest.FieldCondition(
                        key="rag_id",
                        match=rest.MatchValue(value=request.kb_id),
                    )
                ]
            )
            
            # 搜索，获取更多候选结果用于后续过滤
            search_limit = min(request.top_k * 3, 100)  # 获取3倍数量用于后续过滤
            
            err, search_results = self._qdrant.search(
                collection_name=self._rag_collection,
                query_vector=query_vector,
                query_filter=filter_condition,
                limit=search_limit,
                score_threshold=request.similarity_threshold
            )
            
            if err != ErrorCode.SUCCESS:
                self.log.warning(f"_stage1_vector_initial_screening: Qdrant search failed err={err}")
                return []
            
            # 转换结果格式
            initial_results = []
            for result in search_results:
                payload = result.get("payload", {})
                initial_results.append({
                    "id": result.get("id"),
                    "score": result.get("score", 0.0),
                    "segment_id": payload.get("segment_id"),
                    "file_id": payload.get("file_id"),
                    "chunk_index": payload.get("chunk_index"),
                    "text": payload.get("text", ""),
                    "rag_id": payload.get("rag_id"),
                    "binding_user_id": payload.get("binding_user_id")
                })
            
            self.log.info(f"_stage1_vector_initial_screening: found {len(initial_results)} initial results")
            return initial_results
            
        except Exception as e:
            self.log.error(f"_stage1_vector_initial_screening exception: {e}")
            return []
    
    async def _stage2_metadata_filtering(self, initial_results: List[Dict[str, Any]], kb_id: int) -> List[Dict[str, Any]]:
        """阶段2: 元数据过滤"""
        try:
            if not initial_results:
                return []
            
            # 获取有效的segment_id列表
            segment_ids = [result["segment_id"] for result in initial_results if result.get("segment_id")]
            if not segment_ids:
                return []
            
            # 批量查询MySQL中的知识片段状态
            err, segments_info = self.repo.get_segments_by_ids(segment_ids)
            if err != ErrorCode.SUCCESS:
                self.log.warning(f"_stage2_metadata_filtering: get_segments_by_ids failed err={err}")
                return initial_results  # 返回原始结果作为兜底
            
            # 构建segment状态映射
            segment_status_map = {}
            for segment in segments_info:
                segment_id = segment.get("id")
                if segment_id:
                    segment_status_map[segment_id] = {
                        "status": segment.get("status", 0),
                        "is_noise": segment.get("is_noise", False),
                        "is_redundant": segment.get("is_redundant", False)
                    }
            
            # 过滤结果：仅保留启用且非噪音、非冗余的片段
            filtered_results = []
            for result in initial_results:
                segment_id = result.get("segment_id")
                if not segment_id:
                    continue
                
                segment_info = segment_status_map.get(segment_id)
                if segment_info and segment_info["status"] == 2 and not segment_info["is_noise"] and not segment_info["is_redundant"]:
                    filtered_results.append(result)
            
            self.log.info(f"_stage2_metadata_filtering: filtered {len(initial_results)} -> {len(filtered_results)} results")
            return filtered_results
            
        except Exception as e:
            self.log.error(f"_stage2_metadata_filtering exception: {e}")
            return initial_results  # 返回原始结果作为兜底
    
    async def _stage3_semantic_extension(self, query: str, filtered_results: List[Dict[str, Any]]) -> List[str]:
        """阶段3: 语义扩展"""
        try:
            if not self._neo4j_client or not self._llm_client:
                self.log.info("_stage3_semantic_extension: Neo4j or LLM client not available, skipping semantic extension")
                return []
            
            # 提取查询中的关键实体
            entities = await self._extract_entities_from_query(query)
            if not entities:
                return []
            
            # 通过Neo4j图遍历获取相关概念
            related_concepts = await self._get_related_concepts_from_graph(entities)
            
            # 去重并限制数量
            unique_concepts = list(set(related_concepts))[:10]
            
            self.log.info(f"_stage3_semantic_extension: extracted {len(entities)} entities, found {len(unique_concepts)} related concepts")
            return unique_concepts
            
        except Exception as e:
            self.log.error(f"_stage3_semantic_extension exception: {e}")
            return []
    
    async def _extract_entities_from_query(self, query: str) -> List[str]:
        """从查询中提取关键实体"""
        try:
            # 使用LLM提取实体
            prompt = f"""
请从以下查询中提取关键实体（名词、概念、术语），每行一个，不要重复：

查询：{query}

关键实体：
"""
            # 放入线程池执行，避免阻塞
            resp = await asyncio.to_thread(self._llm_client.chat, [{"role": "user", "content": prompt}])
            if not resp:
                return []
            
            # 解析响应
            entities = []
            for line in resp.strip().split('\n'):
                entity = line.strip().strip('-').strip()
                if entity and len(entity) > 1 and len(entity) < 50:
                    entities.append(entity)
            
            return entities[:5]  # 最多返回5个实体
            
        except Exception as e:
            self.log.error(f"_extract_entities_from_query exception: {e}")
            return []
    
    async def _get_related_concepts_from_graph(self, entities: List[str]) -> List[str]:
        """从Neo4j图中获取相关概念"""
        try:
            if not entities:
                return []
            
            related_concepts = []
            
            for entity in entities[:3]:  # 限制处理前3个实体
                # 查找直接相关的节点
                cypher = """
                MATCH (n {name: $entity})
                OPTIONAL MATCH (n)-[r]-(m)
                WHERE m.name IS NOT NULL AND m.name <> $entity
                RETURN DISTINCT m.name as concept, type(r) as relation_type
                LIMIT 10
                """
                # 放入线程池执行，并正确解包结果 (ErrorCode, List)
                err, result = await asyncio.to_thread(
                    self._neo4j_client.execute_query, cypher, {"entity": entity}
                )
                
                if err == ErrorCode.SUCCESS and result:
                    for record in result:
                        concept = record.get("concept")
                        if concept and concept not in related_concepts:
                            related_concepts.append(concept)
            
            return related_concepts
            
        except Exception as e:
            self.log.error(f"_get_related_concepts_from_graph exception: {e}")
            return []
    
    async def _stage4_secondary_retrieval_enhancement(
        self, 
        original_query: str, 
        related_concepts: List[str], 
        kb_id: int, 
        top_k: int
    ) -> List[Dict[str, Any]]:
        """阶段4: 二次检索增强"""
        try:
            if not related_concepts:
                return []
            
            # 构造增强查询
            concept_text = " ".join(related_concepts[:5])
            enhanced_query = f"{original_query} {concept_text}"
            
            # 向量化增强查询
            err, embeddings = self.embed_texts([enhanced_query])
            if err != ErrorCode.SUCCESS or not embeddings:
                self.log.warning(f"_stage4_secondary_retrieval_enhancement: embed failed err={err}")
                return []
            
            enhanced_vector = embeddings[0]
            
            # Qdrant 二次检索
            if not self._qdrant:
                return []
            
            filter_condition = rest.Filter(
                must=[
                    rest.FieldCondition(
                        key="rag_id",
                        match=rest.MatchValue(value=kb_id),
                    )
                ]
            )
            
            err, search_results = self._qdrant.search(
                collection_name=self._rag_collection,
                query_vector=enhanced_vector,
                query_filter=filter_condition,
                limit=top_k,
                score_threshold=0.3  # 降低阈值用于二次检索
            )
            
            if err != ErrorCode.SUCCESS:
                return []
            
            # 转换结果格式
            enhanced_results = []
            for result in search_results:
                payload = result.get("payload", {})
                enhanced_results.append({
                    "id": result.get("id"),
                    "score": result.get("score", 0.0),
                    "segment_id": payload.get("segment_id"),
                    "file_id": payload.get("file_id"),
                    "text": payload.get("text", ""),
                    "chunk_index": payload.get("chunk_index")
                })
            
            self.log.info(f"_stage4_secondary_retrieval_enhancement: found {len(enhanced_results)} enhanced results")
            return enhanced_results
            
        except Exception as e:
            self.log.error(f"_stage4_secondary_retrieval_enhancement exception: {e}")
            return []
    
    async def _stage5_result_mapping_and_merging(
        self, 
        enhanced_results: List[Dict[str, Any]], 
        kb_id: int
    ) -> List[Dict[str, Any]]:
        """阶段5: 结果映射与合并"""
        try:
            if not enhanced_results:
                return []
            
            # 获取segment详细信息
            segment_ids = [result["segment_id"] for result in enhanced_results if result.get("segment_id")]
            if not segment_ids:
                return []
            
            err, segments_info = self.repo.get_segments_by_ids(segment_ids)
            if err != ErrorCode.SUCCESS:
                self.log.warning(f"_stage5_result_mapping_and_merging: get_segments_by_ids failed err={err}")
                return enhanced_results
            
            # 构建segment信息映射
            segment_info_map = {}
            for segment in segments_info:
                segment_id = segment.get("id")
                if segment_id:
                    segment_info_map[segment_id] = segment
            
            # 合并结果
            final_results = []
            seen_content = set()  # 用于去重
            
            for result in enhanced_results:
                segment_id = result.get("segment_id")
                if not segment_id:
                    continue
                
                segment_info = segment_info_map.get(segment_id)
                if not segment_info:
                    continue
                
                # 内容去重
                content = segment_info.get("content", "")
                content_hash = hash(content[:100])  # 使用前100字符的hash进行去重
                if content_hash in seen_content:
                    continue
                seen_content.add(content_hash)
                
                # 构建最终结果
                final_result = {
                    "segment_id": segment_id,
                    "title": segment_info.get("title", ""),
                    "content": content,
                    "summary": segment_info.get("summary", ""),
                    "similarity_score": result.get("score", 0.0),
                    "chunk_index": result.get("chunk_index", 0),
                    "file_id": segment_info.get("file_id"),
                    "recall_count": segment_info.get("recall_count", 0),
                    "created_at": segment_info.get("created_at"),
                    "updated_at": segment_info.get("updated_at")
                }
                
                final_results.append(final_result)
            
            # 按相似度分数排序
            final_results.sort(key=lambda x: x["similarity_score"], reverse=True)
            
            self.log.info(f"_stage5_result_mapping_and_merging: merged {len(enhanced_results)} -> {len(final_results)} final results")
            return final_results
            
        except Exception as e:
            self.log.error(f"_stage5_result_mapping_and_merging exception: {e}")
            return enhanced_results

